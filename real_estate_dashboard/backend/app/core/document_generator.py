"""Document generation utilities for PDF and PowerPoint exports."""

import io
import base64
from typing import Dict, List, Optional, Any, Tuple
from datetime import datetime
from pathlib import Path

# PDF Generation
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.platypus import (
    SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer,
    PageBreak, Image, KeepTogether
)
from reportlab.lib.enums import TA_CENTER, TA_RIGHT, TA_LEFT

# PowerPoint Generation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Chart Generation
import matplotlib
matplotlib.use('Agg')  # Use non-interactive backend
import matplotlib.pyplot as plt
import numpy as np


class PDFGenerator:
    """Generate professional PDF documents with charts and tables."""

    def __init__(self, title: str, author: str = "Real Estate Dashboard"):
        """
        Initialize PDF generator.

        Args:
            title: Document title
            author: Document author
        """
        self.title = title
        self.author = author
        self.styles = getSampleStyleSheet()
        self._setup_custom_styles()

    def _setup_custom_styles(self):
        """Set up custom paragraph styles."""
        # Title style
        self.styles.add(ParagraphStyle(
            name='CustomTitle',
            parent=self.styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#1a237e'),
            spaceAfter=30,
            alignment=TA_CENTER,
            fontName='Helvetica-Bold'
        ))

        # Section heading style
        self.styles.add(ParagraphStyle(
            name='SectionHeading',
            parent=self.styles['Heading2'],
            fontSize=16,
            textColor=colors.HexColor('#283593'),
            spaceAfter=12,
            spaceBefore=12,
            fontName='Helvetica-Bold'
        ))

        # Subsection heading style
        self.styles.add(ParagraphStyle(
            name='SubsectionHeading',
            parent=self.styles['Heading3'],
            fontSize=14,
            textColor=colors.HexColor('#3f51b5'),
            spaceAfter=8,
            spaceBefore=8,
            fontName='Helvetica-Bold'
        ))

        # Body text with better spacing
        self.styles.add(ParagraphStyle(
            name='CustomBody',
            parent=self.styles['Normal'],
            fontSize=11,
            leading=14,
            spaceAfter=6,
            alignment=TA_LEFT
        ))

        # Highlight style
        self.styles.add(ParagraphStyle(
            name='Highlight',
            parent=self.styles['Normal'],
            fontSize=12,
            textColor=colors.HexColor('#1976d2'),
            fontName='Helvetica-Bold',
            spaceAfter=6
        ))

    def generate(
        self,
        sections: List[Dict[str, Any]],
        output_path: Optional[str] = None
    ) -> bytes:
        """
        Generate PDF from sections.

        Args:
            sections: List of section dictionaries with 'type', 'title', 'content', etc.
            output_path: Optional file path to save PDF

        Returns:
            PDF bytes
        """
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(
            buffer,
            pagesize=letter,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=18,
            title=self.title,
            author=self.author
        )

        # Build story
        story = []

        # Add document title
        story.append(Paragraph(self.title, self.styles['CustomTitle']))
        story.append(Spacer(1, 0.2 * inch))

        # Add generation timestamp
        timestamp = datetime.now().strftime("%B %d, %Y")
        story.append(Paragraph(f"Generated: {timestamp}", self.styles['Normal']))
        story.append(Spacer(1, 0.3 * inch))

        # Process each section
        for section in sections:
            section_type = section.get('type', 'text')

            if section_type == 'heading':
                story.append(Paragraph(section['title'], self.styles['SectionHeading']))

            elif section_type == 'subheading':
                story.append(Paragraph(section['title'], self.styles['SubsectionHeading']))

            elif section_type == 'text':
                content = section.get('content', '')
                story.append(Paragraph(content, self.styles['CustomBody']))
                story.append(Spacer(1, 0.1 * inch))

            elif section_type == 'highlight':
                content = section.get('content', '')
                story.append(Paragraph(content, self.styles['Highlight']))
                story.append(Spacer(1, 0.1 * inch))

            elif section_type == 'table':
                table_data = section.get('data', [])
                if table_data:
                    table = self._create_table(table_data, section.get('widths'))
                    story.append(table)
                    story.append(Spacer(1, 0.2 * inch))

            elif section_type == 'chart':
                chart_image = section.get('image')  # base64 or bytes
                if chart_image:
                    img = self._create_image_from_data(chart_image)
                    if img:
                        story.append(img)
                        story.append(Spacer(1, 0.2 * inch))

            elif section_type == 'page_break':
                story.append(PageBreak())

            elif section_type == 'spacer':
                height = section.get('height', 0.2)
                story.append(Spacer(1, height * inch))

        # Build PDF
        doc.build(story)

        # Get PDF bytes
        pdf_bytes = buffer.getvalue()
        buffer.close()

        # Optionally save to file
        if output_path:
            with open(output_path, 'wb') as f:
                f.write(pdf_bytes)

        return pdf_bytes

    def _create_table(self, data: List[List[Any]], widths: Optional[List[float]] = None) -> Table:
        """Create a formatted table."""
        # Convert all data to strings
        table_data = [[str(cell) for cell in row] for row in data]

        # Calculate column widths
        if widths:
            col_widths = [w * inch for w in widths]
        else:
            # Auto-calculate widths
            num_cols = len(table_data[0]) if table_data else 1
            available_width = 6.5 * inch  # Total available width
            col_widths = [available_width / num_cols] * num_cols

        table = Table(table_data, colWidths=col_widths)

        # Apply table style
        table.setStyle(TableStyle([
            # Header row
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#3f51b5')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('TOPPADDING', (0, 0), (-1, 0), 12),

            # Body rows
            ('BACKGROUND', (0, 1), (-1, -1), colors.white),
            ('TEXTCOLOR', (0, 1), (-1, -1), colors.black),
            ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 1), (-1, -1), 10),
            ('TOPPADDING', (0, 1), (-1, -1), 8),
            ('BOTTOMPADDING', (0, 1), (-1, -1), 8),

            # Alternating row colors
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#f5f5f5')]),

            # Grid
            ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ]))

        return table

    def _create_image_from_data(self, image_data: Any, width: float = 6) -> Optional[Image]:
        """Create Image object from base64 or bytes data."""
        try:
            if isinstance(image_data, str):
                # Assume base64
                if image_data.startswith('data:image'):
                    # Strip data URI prefix
                    image_data = image_data.split(',')[1]
                image_bytes = base64.b64decode(image_data)
            else:
                image_bytes = image_data

            # Create BytesIO object
            img_buffer = io.BytesIO(image_bytes)
            img = Image(img_buffer, width=width * inch)
            return img
        except Exception as e:
            print(f"Error creating image: {e}")
            return None


class PowerPointGenerator:
    """Generate professional PowerPoint presentations with charts and tables."""

    def __init__(self, title: str, subtitle: str = ""):
        """
        Initialize PowerPoint generator.

        Args:
            title: Presentation title
            subtitle: Presentation subtitle
        """
        self.title = title
        self.subtitle = subtitle
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

    def add_title_slide(self, title: Optional[str] = None, subtitle: Optional[str] = None):
        """Add title slide."""
        title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_slide_layout)

        title_text = title or self.title
        subtitle_text = subtitle or self.subtitle

        slide.shapes.title.text = title_text
        if subtitle_text:
            slide.placeholders[1].text = subtitle_text

    def add_content_slide(
        self,
        title: str,
        content: List[str] = None,
        table_data: List[List[Any]] = None,
        chart_image: Any = None
    ):
        """
        Add content slide with text, table, or chart.

        Args:
            title: Slide title
            content: List of bullet points
            table_data: 2D list for table
            chart_image: Image data (base64 or bytes)
        """
        # Use blank slide layout
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(0.75)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(26, 35, 126)

        # Add content
        top = Inches(1.5)

        if content:
            # Add bullet points
            text_box = slide.shapes.add_textbox(
                Inches(0.5), top, Inches(9), Inches(5)
            )
            text_frame = text_box.text_frame
            text_frame.word_wrap = True

            for i, bullet in enumerate(content):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = Pt(18)

        elif table_data:
            # Add table
            rows = len(table_data)
            cols = len(table_data[0]) if table_data else 0

            table = slide.shapes.add_table(
                rows, cols,
                Inches(0.5), top,
                Inches(9), Inches(4)
            ).table

            # Populate table
            for i, row in enumerate(table_data):
                for j, cell_value in enumerate(row):
                    cell = table.cell(i, j)
                    cell.text = str(cell_value)

                    # Header row formatting
                    if i == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(63, 81, 181)
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.color.rgb = RGBColor(255, 255, 255)
                                run.font.bold = True
                                run.font.size = Pt(14)
                    else:
                        # Body cells
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(12)

        elif chart_image:
            # Add chart image
            try:
                if isinstance(chart_image, str):
                    # Base64
                    if chart_image.startswith('data:image'):
                        chart_image = chart_image.split(',')[1]
                    image_bytes = base64.b64decode(chart_image)
                else:
                    image_bytes = chart_image

                img_buffer = io.BytesIO(image_bytes)
                slide.shapes.add_picture(
                    img_buffer,
                    Inches(1), top,
                    width=Inches(8)
                )
            except Exception as e:
                print(f"Error adding chart to slide: {e}")

    def save(self, output_path: str):
        """Save presentation to file."""
        self.prs.save(output_path)

    def to_bytes(self) -> bytes:
        """Return presentation as bytes."""
        buffer = io.BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()


class ChartGenerator:
    """Generate charts as images for embedding in documents."""

    @staticmethod
    def create_bar_chart(
        data: Dict[str, float],
        title: str,
        xlabel: str = "",
        ylabel: str = "",
        figsize: Tuple[int, int] = (10, 6)
    ) -> bytes:
        """
        Create bar chart and return as PNG bytes.

        Args:
            data: Dictionary of labels and values
            title: Chart title
            xlabel: X-axis label
            ylabel: Y-axis label
            figsize: Figure size (width, height)

        Returns:
            PNG image bytes
        """
        fig, ax = plt.subplots(figsize=figsize)

        labels = list(data.keys())
        values = list(data.values())

        bars = ax.bar(labels, values, color='#3f51b5')

        # Add value labels on bars
        for bar in bars:
            height = bar.get_height()
            ax.text(
                bar.get_x() + bar.get_width() / 2., height,
                f'${height:,.0f}' if height > 1000 else f'{height:.1f}',
                ha='center', va='bottom'
            )

        ax.set_title(title, fontsize=16, fontweight='bold')
        ax.set_xlabel(xlabel, fontsize=12)
        ax.set_ylabel(ylabel, fontsize=12)
        ax.grid(axis='y', alpha=0.3)

        # Rotate x-axis labels if needed
        if len(labels) > 5:
            plt.xticks(rotation=45, ha='right')

        plt.tight_layout()

        # Save to bytes
        buffer = io.BytesIO()
        plt.savefig(buffer, format='png', dpi=300, bbox_inches='tight')
        buffer.seek(0)
        plt.close(fig)

        return buffer.getvalue()

    @staticmethod
    def create_line_chart(
        data: Dict[str, List[float]],
        x_labels: List[str],
        title: str,
        xlabel: str = "",
        ylabel: str = "",
        figsize: Tuple[int, int] = (10, 6)
    ) -> bytes:
        """
        Create line chart and return as PNG bytes.

        Args:
            data: Dictionary of series names and values
            x_labels: X-axis labels
            title: Chart title
            xlabel: X-axis label
            ylabel: Y-axis label
            figsize: Figure size

        Returns:
            PNG image bytes
        """
        fig, ax = plt.subplots(figsize=figsize)

        for series_name, values in data.items():
            ax.plot(x_labels, values, marker='o', label=series_name, linewidth=2)

        ax.set_title(title, fontsize=16, fontweight='bold')
        ax.set_xlabel(xlabel, fontsize=12)
        ax.set_ylabel(ylabel, fontsize=12)
        ax.legend()
        ax.grid(True, alpha=0.3)

        plt.xticks(rotation=45, ha='right')
        plt.tight_layout()

        # Save to bytes
        buffer = io.BytesIO()
        plt.savefig(buffer, format='png', dpi=300, bbox_inches='tight')
        buffer.seek(0)
        plt.close(fig)

        return buffer.getvalue()

    @staticmethod
    def create_pie_chart(
        data: Dict[str, float],
        title: str,
        figsize: Tuple[int, int] = (8, 8)
    ) -> bytes:
        """
        Create pie chart and return as PNG bytes.

        Args:
            data: Dictionary of labels and values
            title: Chart title
            figsize: Figure size

        Returns:
            PNG image bytes
        """
        fig, ax = plt.subplots(figsize=figsize)

        labels = list(data.keys())
        values = list(data.values())

        colors = plt.cm.Set3(range(len(labels)))

        wedges, texts, autotexts = ax.pie(
            values,
            labels=labels,
            autopct='%1.1f%%',
            startangle=90,
            colors=colors
        )

        # Beautify text
        for text in texts:
            text.set_fontsize(12)
        for autotext in autotexts:
            autotext.set_color('white')
            autotext.set_fontweight('bold')
            autotext.set_fontsize(10)

        ax.set_title(title, fontsize=16, fontweight='bold')

        plt.tight_layout()

        # Save to bytes
        buffer = io.BytesIO()
        plt.savefig(buffer, format='png', dpi=300, bbox_inches='tight')
        buffer.seek(0)
        plt.close(fig)

        return buffer.getvalue()
